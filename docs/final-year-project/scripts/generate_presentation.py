#!/usr/bin/env python3
"""Generate CamTraffic final defense presentation (15 slides). Task 407."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUT = Path(__file__).resolve().parent.parent / "CAMTRAFFIC-FINAL-PRESENTATION.pptx"

TITLE_COLOR = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0x25, 0x63, 0xEB)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2))
        tf2 = box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
    return slide


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8))
    tb.text_frame.paragraphs[0].text = title
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.bold = True
    tb.text_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.space_after = Pt(8)
        if line.startswith("•") or line.startswith("-"):
            p.level = 0
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "CamTraffic",
        "AI-Based Traffic Sign Detection and Traffic Law Enforcement System\nCambodia · Final Year Project · 2026\n[Your Name] · [University]",
    )

    add_content_slide(
        prs,
        "Problem Statement",
        [
            "• Officers cannot monitor all intersections continuously",
            "• Violations lack standardized photographic evidence",
            "• Drivers have limited transparency into fines and appeals",
            "• Detection, fines, and appeals operate in separate processes",
            "",
            "CamTraffic connects AI detection → violations → fines → appeals",
        ],
    )

    add_content_slide(
        prs,
        "Project Objectives",
        [
            "O1 Detect Cambodian traffic signs (YOLO) — ✅ mAP@50 = 0.908",
            "O2 License plate OCR integration — ✅ Pipeline live",
            "O3 Violation → fine → appeal workflow — ✅ Full REST API",
            "O4 Admin + officer + driver portals — ✅ React SPAs",
            "O5 Production deployment — ✅ Docker 8-service stack",
        ],
    )

    add_content_slide(
        prs,
        "System Overview",
        [
            "Admin portal (:5174) — users, cameras, reports, audit",
            "User portal (:5173) — Officer + Driver tabs",
            "Django REST API + PostgreSQL",
            "YOLO11n + EasyOCR embedded pipeline",
            "",
            "Three stakeholders · One enforcement platform",
        ],
    )

    add_content_slide(
        prs,
        "Architecture (Task 403)",
        [
            "Modular monolith — Django REST + embedded AI",
            "",
            "Client: React 19 + Vite (admin + user)",
            "API: Django 5 + Gunicorn + DRF",
            "AI: YOLO11n + EasyOCR (+ optional ai-worker)",
            "Data: PostgreSQL 16, Redis 7",
            "Edge: Nginx + Let's Encrypt (production)",
            "",
            "Diagram: docs/final-year-project/diagrams/DEPLOYMENT-DIAGRAM.md",
        ],
    )

    add_content_slide(
        prs,
        "Use Cases & Roles",
        [
            "Admin — users, RBAC, cameras, backup, audit logs",
            "Traffic Police — AI detect, violations, fines, appeals review",
            "Driver — view fines, pay (demo), submit appeals, vehicles",
            "",
            "Figure: USE-CASE-DIAGRAM.md",
        ],
    )

    add_content_slide(
        prs,
        "AI Pipeline (Task 404)",
        [
            "1. Image upload or webcam frame",
            "2. YOLO11n sign detection (10 classes, 640px)",
            "3. Vehicle bbox + EasyOCR plate read",
            "4. ViolationRule engine (sign_class_key match)",
            "5. AIDetectionLog → Violation → Fine → Notification",
            "",
            "API: POST /api/ai/detect/ · Weights: ai/weights/best_v2.pt",
        ],
    )

    add_content_slide(
        prs,
        "Dataset (Task 405)",
        [
            "Full sign taxonomy: 31 classes (CVAT registry)",
            "Road footage collection target: 8,848 frames",
            "Traffic sign images collected: 2,840",
            "Grand total images (all categories): 3,561",
            "",
            "Production model: 10-class subset (dataset_10/)",
            "Classes: NO_ENTRY, M_STOP, speed limits, turn restrictions, etc.",
        ],
    )

    add_content_slide(
        prs,
        "Model Results (Task 406)",
        [
            "YOLO11n — 10-class production model (best_v2.pt)",
            "",
            "mAP@50: 0.908 (target ≥ 0.85) ✅",
            "mAP@50-95: 0.796",
            "Precision: 0.960 · Recall: 0.196",
            "CPU FPS @ 640px: ~20",
            "",
            "Insert: PR_curve.png · confusion_matrix.png",
            "OCR CER: 2.40 — officer confirmation required",
        ],
    )

    add_content_slide(
        prs,
        "Live Demo (7 Scenes)",
        [
            "1. Admin dashboard KPIs",
            "2. Live camera grid",
            "3. AI detection (upload/webcam)",
            "4. Violation auto-create + notification",
            "5. Officer confirm + issue fine",
            "6. Driver portal + payment demo",
            "7. PDF/Excel report export",
            "",
            "Script: DEMO-SCRIPT.md (~12 minutes)",
        ],
    )

    add_content_slide(
        prs,
        "Testing & UAT",
        [
            "Backend API + RBAC tests — PASS",
            "Security (SQL injection, upload MIME) — PASS",
            "Playwright E2E — 4/4 PASS",
            "UAT (admin, officer, driver) — PASS",
            "Health endpoint p95 < 250 ms — PASS",
            "",
            "Evidence: UAT-REPORT.md · PERFORMANCE-EVALUATION.md",
        ],
    )

    add_content_slide(
        prs,
        "Production Deployment",
        [
            "Docker Compose — 8 services:",
            "nginx · backend · ai-worker · celery-worker · celery-beat",
            "postgres · redis · certbot",
            "",
            "npm run docker:prod:up",
            "CI: .github/workflows/ci.yml",
        ],
    )

    add_content_slide(
        prs,
        "Key Achievements",
        [
            "446 / 540 checklist tasks completed",
            "~120 API endpoints · 16 Django apps",
            "7 thesis chapters · 40+ documentation files",
            "",
            "Hybrid: AI perception + expert-system rules + full-stack IT",
        ],
    )

    add_content_slide(
        prs,
        "Limitations & Future Work",
        [
            "Limitations:",
            "• 10-class vs 31-class catalog",
            "• OCR needs Cambodia-specific training",
            "• Web-only · demo payment",
            "",
            "Future: expand dataset, GPU RTSP, mobile app, police pilot",
        ],
    )

    add_title_slide(
        prs,
        "Thank You",
        "Questions?\n\ngithub.com/SareachGenZ/CamTraffic\nDEFENSE-PREPARATION.md",
    )

    prs.save(OUT)
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
